from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_core.language_models import LLM
from typing import Optional, List
import streamlit as st 

from multimodal import get_response


import pickle
import os


import re
from logger import log_execution

# Document loaders
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,
    CSVLoader,
    UnstructuredPowerPointLoader
)

from langchain_text_splitters import RecursiveCharacterTextSplitter

from pathlib import Path

from rank_bm25 import BM25Okapi

from pptx import Presentation
from PIL import Image
import io
from langchain_core.documents import Document
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Stores documents globally
vectorless_docs = []
bm25_model = None
CHUNKS_FILE = "vectorless_chunks.pkl"
# Load chunks if file exists or Load saved chunks
if os.path.exists(CHUNKS_FILE):

    with open(CHUNKS_FILE, "rb") as f:

        data = pickle.load(f)
        
        if isinstance(data, list):   #list format storage
            vectorless_docs = data
            bm25_model = None
  
        else:         #dictionary format storage
            vectorless_docs = data["docs"]
            bm25_model = data["bm25"]

        vectorless_docs = data["docs"]
        bm25_model = data["bm25"]

        print("Loaded vectorless docs:", len(vectorless_docs))

# Custom LLM
class CustomLLM(LLM):

    def _call(self, prompt: str, stop: Optional[List[str]] = None) -> str:
        return get_response(prompt)

    @property
    def _identifying_params(self):
        return {}

    @property
    def _llm_type(self):
        return "custom_groq"


@log_execution("llm_generation")
def generate_response(llm, prompt):
    return llm.invoke(prompt)



@log_execution("reranking")
def rerank_documents(query, docs, top_k=5):

    if not docs:
        return docs

    query_tokens = re.findall(r"\w+", query.lower())

    scored_docs = []

    for doc in docs:

        text = doc.page_content.lower()

        score = sum(
            1 for token in query_tokens 
            if token in text
        )

        if score >= 3:   
            scored_docs.append((score, doc))

    ranked_docs = sorted(
        scored_docs,
        key=lambda x: x[0],
        reverse=True
    )

    return [doc for _, doc in ranked_docs[:top_k]]


@log_execution("image_llm_generation")
def generate_image_response(prompt, image):
    return get_response(prompt, image)


@log_execution("post_processing")
def post_process_response(response):
    return response.strip()


# Load documents
def load_documents(file_path):

    file_path = str(file_path).lower()

    if file_path.endswith(".pdf"):
        loader = PyPDFLoader(file_path)

    elif file_path.endswith(".docx"):
        loader = Docx2txtLoader(file_path)

    elif file_path.endswith(".txt"):
        loader = TextLoader(file_path)

    elif file_path.endswith(".csv"):
        loader = CSVLoader(file_path)

    elif file_path.endswith(".pptx"):
        loader = UnstructuredPowerPointLoader(file_path)

    else:
        raise ValueError(f"Unsupported file format: {file_path}")

    return loader.load()


def extract_images_from_ppt(ppt_path):

    images = []

    prs = Presentation(ppt_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

                image = shape.image
                image_bytes = image.blob

                images.append(image_bytes)

    return images

# Load vectorless docs
def load_vectorless_docs(pdf_paths):

    global vectorless_docs
    global bm25_model

    all_documents = []

    for file_path in pdf_paths:

        file_path = Path(file_path)

        documents = load_documents(file_path)

        if documents:
            for doc in documents:
                doc.metadata = {"source": file_path.name}
                all_documents.append(doc)

        # Extract images from PPT
        if file_path.suffix == ".pptx":

            images = extract_images_from_ppt(file_path)

            if images:
                progress = st.progress(0)

                for i, img in enumerate(images):
                    try:
                        caption = get_response(
                            "Describe this image in detail",
                            img
                        )

                        all_documents.append(
                            Document(
                                page_content=caption,
                                metadata={"source": file_path.name}
                            )
                        )

                        progress.progress((i + 1) / len(images))

                    except Exception as e:
                        print(f"Image extraction error: {e}")

    # Split documents
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100
    )

    new_docs = text_splitter.split_documents(all_documents)

    vectorless_docs.extend(new_docs)

    # rebuild bm25
    tokenized_docs = [
        doc.page_content.lower().split()
        for doc in vectorless_docs
]

    bm25_model = BM25Okapi(tokenized_docs)

    with open(CHUNKS_FILE, "wb") as f:
        pickle.dump({
            "docs": vectorless_docs,
            "bm25": bm25_model
        }, f)

    st.write("Chunks stored:", len(vectorless_docs))


def filter_best_document(docs):

    if not docs:
        return docs

    doc_scores = {}

    for doc in docs:
        source = doc.metadata.get("source", "Unknown")

        if source not in doc_scores:
            doc_scores[source] = 0

        doc_scores[source] += 1

    if not doc_scores:
        return docs

    best_source = max(doc_scores, key=doc_scores.get)

    filtered_docs = [
        doc for doc in docs
        if doc.metadata.get("source") == best_source
    ]

    return filtered_docs


# Load RAG
def load_rag(chat_session_id=None):

    print("Loading Vector-less RAG...")

    llm = CustomLLM()

    def format_docs(docs):
        return "\n\n".join(doc.page_content for doc in docs)

    def vectorless_rag(query):

        docs = hybrid_retrieval(query)

        docs = rerank_documents(query, docs)

        docs = filter_best_document(docs)

        if not docs:
            response = generate_response(llm, query)
            return post_process_response(response)

        context = format_docs(docs)

        
        prompt = create_rag_prompt(context, query)

        response = generate_response(llm, prompt)

        response = post_process_response(response)

        best_source = docs[0].metadata.get("source", "Unknown")

        response += "\n\n**Sources:**\n"
        response += f"- {best_source}\n"

        return response
    return vectorless_rag


@log_execution("prompt_creation")
def create_prompt(context, query):
    return f"""
You are a helpful AI assistant.

Use the provided context to answer the question completely.
If the answer exists in context, return full information.

If context is not relevant, answer using general knowledge.

Context:
{context}

Question:
{query}

Answer:
"""


# Save docs when uploaded
def save_pdfs_to_vector_db(pdf_paths, chat_session_id):

    load_vectorless_docs(pdf_paths)

def get_persistent_retriever(chat_session_id):
    return None

@log_execution("bm25_retrieval")
def bm25_retrieval(query):

    global vectorless_docs
    global bm25_model

    if not vectorless_docs:
        return []

    query_tokens = preprocess_query(query)

    scores = bm25_model.get_scores(query_tokens)
    scores = list(scores)

    scored_docs = list(zip(scores, vectorless_docs))

    scored_docs = sorted(
        scored_docs,
        key=lambda x: x[0],
        reverse=True
    )

    
    top_score = scored_docs[0][0] if scored_docs else 0

    filtered_docs = [
        doc for score, doc in scored_docs
        if score >= top_score * 0.6   
    ]

    return filtered_docs[:5]


def preprocess_query(query):

    stopwords = {
        "what","is","are","the","of","in","for",
        "does","do","a","an","and","to"
    }

    tokens = re.findall(r"\w+", query.lower())

    filtered = [
        token for token in tokens
        if token not in stopwords
    ]

    return filtered




def keyword_retrieval(query):

    global vectorless_docs

    query_tokens = query.lower().split()

    matched_docs = []

    for doc in vectorless_docs:
        text = doc.page_content.lower()

        if any(token in text for token in query_tokens):
            matched_docs.append(doc)

    return matched_docs[:5]


def hybrid_retrieval(query):

    top_k = dynamic_top_k(query)

    bm25_docs = bm25_retrieval(query)

    keyword_docs = keyword_retrieval(query)

    combined_docs = bm25_docs + keyword_docs

    unique_docs = []
    seen = set()

    for doc in combined_docs:
        content = doc.page_content

        if content not in seen:
            unique_docs.append(doc)
            seen.add(content)

    return unique_docs[:top_k]

def dynamic_top_k(query):

    query = query.lower()

    if any(word in query for word in ["list", "keywords", "types", "advantages"]):
        return 6

    if any(word in query for word in ["what is", "define", "meaning"]):
        return 2

    if any(word in query for word in ["difference", "compare"]):
        return 4

    return 5


def create_rag_prompt(context, query):
    return f"""
You are a helpful AI assistant.

Answer ONLY from the provided context.
Do NOT add extra information.

Context:
{context}

Question:
{query}

Answer:
"""
